from pathlib import Path
from typing import List
from langchain_core.documents import Document
from backend.app.utils.logger import logger

def load_pptx(file_path: str) -> List[Document]:
    """Loads a PPTX document and extracts slides as LangChain Documents."""
    logger.info(f"Starting PPTX loading for file: {file_path}")
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        docs = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_runs = []
            
            # Extract text from shapes (text boxes, tables, etc.)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.strip():
                        slide_text_runs.append(shape.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text_runs.append(cell.text.strip())
                    
            slide_content = "\n".join(slide_text_runs)
            
            # Skip empty slides
            if not slide_content.strip():
                continue
                
            doc = Document(
                page_content=slide_content,
                metadata={
                    "file_name": path.name,
                    "file_type": "PPTX",
                    "page_number": slide_idx + 1,  # 1-indexed slide number
                    "source": str(path)
                }
            )
            docs.append(doc)
            
        logger.info(f"Successfully loaded PPTX. Slides: {len(docs)}")
        return docs
    except ImportError:
        logger.error("python-pptx package is not installed. PowerPoint loading is disabled.")
        raise ImportError("python-pptx is required to parse PPTX files.")
    except Exception as e:
        logger.error(f"Error loading PPTX file {file_path}: {e}", exc_info=True)
        raise e
